#!/usr/bin/env python3
"""生成「改稿对比」风格 PPT：横屏 16:9 · 左右两块 · 牛油果绿外框 · 波浪分割 · 黑体。"""
from __future__ import annotations

import math
import random
import shutil
from pathlib import Path
from typing import Sequence

import fitz  # pymupdf
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
ASSETS_SRC = Path(r"C:\Users\Administrator\.cursor\projects\d-Resume-project\assets")
PPT_DIR = ROOT / "public" / "demo" / "ppt"
SLIDES_DIR = PPT_DIR / "slides-xhs"
OUT_PPTX = ROOT / "public" / "demo" / "Linkola-简历对比.pptx"
DESKTOP_PPTX = Path(r"c:\Users\Administrator\Desktop\Linkola-resume-comparison.pptx")
PAPER_PDF = Path(r"c:\Users\Administrator\Desktop\传统纸质例子.pdf")
UPLOADED_PAPER = ASSETS_SRC / (
    "c__Users_Administrator_AppData_Roaming_Cursor_User_workspaceStorage_"
    "58b12e0ea291394258a5ef5f79e972bc_images_image-3e8b7232-3212-478e-83a8-c6f52f8ee42b.png"
)

# 横屏 16:9
W, H = 1920, 1080
PAD = 52
COL_GAP = 28

AVO = (90, 107, 72)
CREAM = (250, 247, 240)
INK = (30, 30, 30)
MUTED = (92, 92, 92)
DARK_PANEL = (46, 42, 36)
WAVE_COLORS = [(168, 184, 138), (123, 140, 92), (90, 107, 72)]
WHITE = (255, 255, 255)

WEB_MAP = {
    "home": "0bc1a63c789204abc8c2be8336a1bd9e",
    "resume": "2a907a9a9e9a12a82f6f1f78a0984cda",
    "scroll": "f5e03e01b3942f55d74426339e86e775",
    "portfolio": "8069fd268098dfd2d27d506d73bca57a",
    "portfolio2": "94a6154fb7a59630d4c9b767fbbbb0a3",
    "home2": "7d6da8e09ef19b878c9d47069527316d",
    "work_detail": "e415dfc469d2e51ffc58541360283166",
    "edu_detail": "383c54bba7da8e82941e4b8aed431bad",
    "hr_spread": "0920a2eaa3e759899e8a94e46ca62f03",
}

SlideSpec = dict


def ensure_assets() -> dict[str, Path]:
    PPT_DIR.mkdir(parents=True, exist_ok=True)
    paths: dict[str, Path] = {}

    if PAPER_PDF.exists():
        shutil.copy2(PAPER_PDF, PPT_DIR / "paper.pdf")
    paper_pdf = PPT_DIR / "paper.pdf"
    if paper_pdf.exists():
        doc = fitz.open(paper_pdf)
        for i, page in enumerate(doc):
            out = PPT_DIR / f"paper-page-{i + 1}.png"
            if not out.exists() or out.stat().st_mtime < paper_pdf.stat().st_mtime:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                pix.save(out.as_posix())
            paths[f"paper_{i + 1}"] = out
        doc.close()

    if UPLOADED_PAPER.exists():
        dest = PPT_DIR / "paper-upload.png"
        shutil.copy2(UPLOADED_PAPER, dest)
        paths["paper_upload"] = dest

    for key, token in WEB_MAP.items():
        dest = PPT_DIR / f"{key}.png"
        if dest.exists():
            paths[key] = dest
            continue
        match = next((p for p in ASSETS_SRC.glob("*.png") if token in p.name), None)
        if match:
            shutil.copy2(match, dest)
            paths[key] = dest

    for key in WEB_MAP:
        dest = PPT_DIR / f"{key}.png"
        if dest.exists():
            paths.setdefault(key, dest)

    return paths


def make_texture(path: Path) -> Image.Image:
    base = AVO
    img = Image.new("RGB", (W, H), base)
    draw = ImageDraw.Draw(img)
    random.seed(42)
    for x in range(0, W, 6):
        tone = random.randint(-18, 18)
        c = tuple(max(0, min(255, base[i] + tone - (i * 3))) for i in range(3))
        draw.line([(x, 0), (x, H)], fill=c, width=2)
    for y in range(0, H, 40):
        draw.line([(0, y), (W, y)], fill=(base[0] - 8, base[1] - 6, base[2] - 10), width=1)
    img = img.filter(ImageFilter.GaussianBlur(radius=1.2))
    overlay = Image.new("RGBA", (W, H), (45, 55, 38, 35))
    img = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")
    img.save(path, quality=92)
    return img


def load_font(size: int, bold: bool = True) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """统一黑体：SimHei / 微软雅黑 Bold。"""
    if bold:
        candidates = [
            r"C:\Windows\Fonts\simhei.ttf",
            r"C:\Windows\Fonts\msyhbd.ttc",
            r"C:\Windows\Fonts\msyh.ttc",
        ]
    else:
        candidates = [
            r"C:\Windows\Fonts\msyh.ttc",
            r"C:\Windows\Fonts\simhei.ttf",
        ]
    for fp in candidates:
        if Path(fp).exists():
            try:
                return ImageFont.truetype(fp, size)
            except OSError:
                continue
    return ImageFont.load_default()


def crop_box(img: Image.Image, box: tuple[float, float, float, float]) -> Image.Image:
    l, t, r, b = box
    w, h = img.size
    return img.crop((int(l * w), int(t * h), int(r * w), int(b * h)))


def load_crop(path: Path | None, box: tuple[float, float, float, float] | None = None) -> Image.Image | None:
    if not path or not path.exists():
        return None
    img = Image.open(path).convert("RGB")
    if box:
        img = crop_box(img, box)
    return img


def fit_contain(img: Image.Image, tw: int, th: int, bg: tuple[int, int, int] = CREAM) -> Image.Image:
    canvas = Image.new("RGB", (tw, th), bg)
    iw, ih = img.size
    scale = min(tw / iw, th / ih)
    nw, nh = max(1, int(iw * scale)), max(1, int(ih * scale))
    resized = img.resize((nw, nh), Image.Resampling.LANCZOS)
    canvas.paste(resized, ((tw - nw) // 2, (th - nh) // 2))
    return canvas


def draw_waves_vertical(
    draw: ImageDraw.ImageDraw,
    x0: int,
    y0: int,
    height: int,
    amplitude: int = 14,
) -> int:
    """在右侧面板左缘绘制竖向三层波浪，返回占用宽度。"""
    total = 0
    for i, color in enumerate(WAVE_COLORS):
        freq = 0.018 + i * 0.003
        phase = i * 1.1
        layer_amp = amplitude - i * 3
        base_x = x0 + i * 8
        pts: list[tuple[int, int]] = [(base_x + layer_amp, y0)]
        for y in range(y0, y0 + height + 1, 3):
            x = base_x + layer_amp * math.sin((y - y0) * freq + phase)
            pts.append((int(x), y))
        edge = base_x + layer_amp + 20
        pts.extend([(edge, y0 + height), (base_x, y0 + height)])
        draw.polygon(pts, fill=color)
        total = max(total, edge - x0)
    return total + 4


def draw_text_block(
    draw: ImageDraw.ImageDraw,
    lines: Sequence[str],
    x: int,
    y: int,
    font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
    fill: tuple[int, int, int],
    line_gap: int = 8,
) -> int:
    cy = y
    for line in lines:
        if not line:
            cy += font.size // 2
            continue
        draw.text((x, cy), line, font=font, fill=fill)
        cy += font.size + line_gap
    return cy


def resolve_img(paths: dict[str, Path], key: str | None, fallback: Path | None = None) -> Path | None:
    if key and key in paths:
        return paths[key]
    if key:
        p = PPT_DIR / f"{key}.png"
        if p.exists():
            return p
    return fallback


def render_panel_column(
    slide: Image.Image,
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    col_w: int,
    panel_h: int,
    label: str,
    label_font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
    img: Image.Image | None,
    caption: Sequence[str],
    cap_font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
    *,
    dark: bool = False,
    text_only_lines: Sequence[str] | None = None,
    cap_font_white: ImageFont.FreeTypeFont | ImageFont.ImageFont | None = None,
) -> None:
    draw.text((x + 6, y), label, font=label_font, fill=WHITE)
    y += label_font.size + 10

    inset = 12
    box = (x, y, x + col_w, y + panel_h)
    bg = DARK_PANEL if dark else WHITE
    outline = (70, 64, 56) if dark else (220, 214, 204)
    draw.rounded_rectangle(box, radius=12, fill=bg, outline=outline, width=2)

    inner_w = col_w - 2 * inset
    inner_h = panel_h - 2 * inset
    cap_lines = list(caption)
    cap_h = len(cap_lines) * (cap_font.size + 6) + 8 if cap_lines else 0
    img_h = inner_h - cap_h

    content_x = x + inset
    content_y = y + inset

    if dark:
        wave_w = draw_waves_vertical(draw, content_x, content_y, inner_h - cap_h, amplitude=11)
        content_x += wave_w + 4
        inner_w -= wave_w + 4

    if text_only_lines and cap_font_white:
        draw_text_block(draw, text_only_lines, content_x + 8, content_y + 16, cap_font_white, WHITE, 12)
    elif img and img_h > 60:
        fitted = fit_contain(img, inner_w, img_h, bg)
        slide.paste(fitted, (content_x, content_y))

    if cap_lines and not text_only_lines:
        cap_y = y + panel_h - cap_h - inset + 4
        cap_fill = (210, 215, 200) if dark else MUTED
        draw_text_block(draw, cap_lines, x + inset + 4, cap_y, cap_font, cap_fill, 4)


def render_slide(texture: Image.Image, spec: SlideSpec, paths: dict[str, Path]) -> Image.Image:
    slide = texture.copy()
    draw = ImageDraw.Draw(slide)

    inner_w = W - 2 * PAD
    col_w = (inner_w - COL_GAP) // 2

    title_font = load_font(38, bold=True)
    sub_font = load_font(20, bold=False)
    label_font = load_font(24, bold=True)
    cap_font = load_font(16, bold=False)
    cap_bold = load_font(18, bold=True)
    mid_font = load_font(17, bold=False)
    foot_font = load_font(15, bold=False)

    y = PAD
    title = spec.get("title", "")
    subtitle = spec.get("subtitle", "")

    if title:
        tw = draw.textlength(title, font=title_font)
        draw.text(((W - tw) / 2, y), title, font=title_font, fill=CREAM)
        y += title_font.size + 8
    if subtitle:
        tw = draw.textlength(subtitle, font=sub_font)
        draw.text(((W - tw) / 2, y), subtitle, font=sub_font, fill=(228, 234, 220))
        y += sub_font.size + 14

    if spec.get("mid_lines"):
        box_h = len(spec["mid_lines"]) * 24 + 22
        draw.rounded_rectangle(
            (PAD, y, PAD + inner_w, y + box_h),
            radius=10,
            fill=CREAM,
            outline=(216, 208, 192),
            width=2,
        )
        draw_text_block(draw, spec["mid_lines"], PAD + 20, y + 10, mid_font, INK, 5)
        y += box_h + 16

    footer = spec.get("footer", "")
    foot_h = foot_font.size + 16 if footer else 0
    panel_h = spec.get("panel_h", H - y - PAD - foot_h - 72)

    left_x = PAD
    right_x = PAD + col_w + COL_GAP

    left_label = spec.get("left_label") or spec.get("top_label", "纸质 PDF · 传统版")
    right_label = spec.get("right_label") or spec.get("bottom_label", "在线网站 · 改后版")

    left_key, left_crop = spec.get("left_img") or spec.get("top_img", ("paper_1", (0, 0, 1, 0.42)))
    left_path = resolve_img(paths, left_key, PPT_DIR / "paper-page-1.png")
    left_img = load_crop(left_path, left_crop)
    left_cap = spec.get("left_caption") or spec.get("top_caption", [])

    render_panel_column(
        slide, draw, left_x, y, col_w, panel_h,
        left_label, label_font, left_img, left_cap, cap_font,
        dark=False,
    )

    right_img_spec = spec.get("right_img") or spec.get("bottom_img")
    right_img = None
    if right_img_spec is not None:
        rk, rc = right_img_spec
        right_img = load_crop(resolve_img(paths, rk, PPT_DIR / "home.png"), rc)

    right_cap = spec.get("right_caption") or spec.get("bottom_caption", [])
    text_only = spec.get("right_text_only") or spec.get("bottom_text_only", False)

    render_panel_column(
        slide, draw, right_x, y, col_w, panel_h,
        right_label, label_font,
        right_img if not text_only else None,
        [] if text_only else right_cap,
        cap_font,
        dark=True,
        text_only_lines=right_cap if text_only else None,
        cap_font_white=cap_bold if text_only else None,
    )

    if footer:
        fw = draw.textlength(footer, font=foot_font)
        draw.text(
            ((W - fw) / 2, H - PAD - foot_font.size),
            footer,
            font=foot_font,
            fill=(210, 218, 198),
        )

    return slide


def build_slide_specs() -> list[SlideSpec]:
    return [
        {
            "title": "同一份简历，两种形态",
            "subtitle": "虚构 demo · XXX · 在线链接 vs 传统 PDF",
            "left_img": ("paper_upload", (0.02, 0.0, 0.98, 0.38)),
            "right_img": ("home", (0.0, 0.0, 1.0, 0.62)),
            "left_caption": ["一页纸扫完 · 排版靠压缩"],
            "right_caption": ["分区浏览 · 亮点卡片 · 扫码直达"],
            "footer": "求意见：你更想收到链接还是 PDF？评论区 A / B / C",
        },
        {
            "title": "先公平对比：内容是一样的",
            "subtitle": "姓名、岗位、经历、教育、技能 — 同一套 demo 数据",
            "mid_lines": [
                "· 不是「网站更满、PDF 更空」，而是同一人物 XXX 的两套呈现",
                "· 差异在形态与阅读方式，不在信息真假",
            ],
            "panel_h": 620,
            "left_img": ("paper_1", (0.05, 0.0, 0.95, 0.35)),
            "right_img": ("home", (0.0, 0.05, 1.0, 0.58)),
            "left_caption": ["PDF 首屏：抬头 + 求职意向"],
            "right_caption": ["网站首屏：同信息，更有层次"],
            "footer": "以下逐条对比：谁更适合 HR / 求职者",
        },
        {
            "title": "差异 ① 第一印象",
            "subtitle": "PDF：一页扫完 · 网站：分区浏览 + 视觉层次",
            "left_img": ("paper_upload", (0.0, 0.0, 1.0, 0.32)),
            "right_img": ("home2", (0.0, 0.0, 1.0, 0.55)),
            "left_caption": ["文字密度高 · 证件照占位 · 模板感"],
            "right_caption": ["亮点拆卡片 · 技能标签 · 二维码联系"],
            "footer": "PDF 卖点：打印友好 · 网站卖点：视觉层次",
        },
        {
            "title": "差异 ② 成果能不能「点开看」",
            "subtitle": "PDF 只能写 · 网站能展示图 / 视频 / 代码",
            "left_img": ("paper_1", (0.0, 0.28, 1.0, 0.58)),
            "right_img": ("work_detail", (0.0, 0.0, 1.0, 0.85)),
            "left_caption": ["工作经历：纯文字描述成果"],
            "right_caption": ["代表项目可预览 · 量化指标更直观"],
            "footer": "前端 / 设计岗：链接里直接看 demo，比「详见附件」少一步",
        },
        {
            "title": "差异 ③ HR 专用浏览",
            "subtitle": "网站独有：HR 摊开模式 · 快速扫关键成果",
            "left_img": ("paper_1", (0.0, 0.32, 1.0, 0.62)),
            "right_img": ("hr_spread", (0.0, 0.0, 1.0, 0.9)),
            "left_caption": ["PDF：自己用眼睛找重点"],
            "right_caption": ["摊开全部成果 + 项目 · 省下载解压"],
            "footer": "HR 时间紧：链接打开即可扫",
        },
        {
            "title": "差异 ④ 视觉型作品",
            "subtitle": "摄影 / 设计 / 视频类 — 网站栅格 vs PDF 文字",
            "left_img": ("paper_1", (0.0, 0.55, 1.0, 0.78)),
            "right_img": ("portfolio", (0.0, 0.0, 1.0, 0.75)),
            "left_caption": ["PDF：技能/项目只能文字列举"],
            "right_caption": ["作品集栅格 + 封面 + 一键跳转"],
            "footer": "适合作品驱动型岗位",
        },
        {
            "title": "差异 ⑤ 结构化深度",
            "subtitle": "岗位匹配项 · 教育/项目详情页",
            "left_img": ("paper_2", (0.0, 0.0, 1.0, 0.55)),
            "right_img": ("edu_detail", (0.0, 0.0, 1.0, 0.85)),
            "left_caption": ["PDF 第二页：继续文字堆叠"],
            "right_caption": ["详情弹层 · 校园成果 · 分层点击"],
            "footer": "网站：信息有层次 · PDF：超过一页就要翻页",
        },
        {
            "title": "别忽视 PDF：它仍然很重要",
            "subtitle": "系统兼容 · 附件习惯 · 归档打印",
            "right_label": "结论 · 不是二选一",
            "left_img": ("paper_2", (0.05, 0.0, 0.95, 0.95)),
            "right_text_only": True,
            "left_caption": ["招聘网投递 · 邮件附件 — 默认格式"],
            "right_caption": [
                "· 一页纸扫完，不依赖网络",
                "· 打印、存档、转发 — 流程成熟",
                "结论：链接展示力 + PDF 兼容性",
            ],
            "footer": "不是二选一，而是各发各的",
        },
        {
            "title": "你更投哪个？",
            "subtitle": "评论区扣字母 · 帮 demo 定方向",
            "panel_h": 640,
            "left_img": ("paper_upload", (0.0, 0.0, 1.0, 0.45)),
            "right_img": ("home", (0.0, 0.0, 1.0, 0.5)),
            "left_caption": ["A · 只要 PDF（兼容习惯）"],
            "right_caption": ["B · 只要链接（展示力）· C · 两个都要"],
            "footer": "本对比为产品演示 · 人物 XXX 虚构 · 非真实求职",
        },
    ]


def main() -> None:
    paths = ensure_assets()
    texture_path = PPT_DIR / "bg-texture-landscape.png"
    texture = make_texture(texture_path)

    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    specs = build_slide_specs()
    slide_paths: list[Path] = []

    for i, spec in enumerate(specs, start=1):
        img = render_slide(texture, spec, paths)
        out = SLIDES_DIR / f"slide-{i:02d}.png"
        img.save(out, quality=95)
        slide_paths.append(out)
        print(f"  slide {i}: {out.name}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for sp in slide_paths:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(sp), 0, 0, width=prs.slide_width, height=prs.slide_height)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"Saved PPT: {OUT_PPTX}")

    try:
        shutil.copy2(OUT_PPTX, DESKTOP_PPTX)
        print(f"Copied to desktop: {DESKTOP_PPTX}")
    except OSError as e:
        print(f"Desktop copy skipped: {e}")


if __name__ == "__main__":
    main()
